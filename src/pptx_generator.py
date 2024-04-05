from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

def initialize_presentation():
    """
    Initializes a PowerPoint presentation object.
    """
    prs = Presentation()
    return prs

def add_textbox(slide, text, left, top, width, height, position, font_size):
    """
    Adds a text box to a slide with specified text, position, size, alignment, and font size.
    """
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = position

    for run in p.runs:
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.size = Pt(font_size)

def add_verse_slide(prs, book_name_kor, book_name_eng, chapter, verse, text_kor, text_eng):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    verse_reference = f"{book_name_kor} {book_name_eng} {chapter}:{verse}"
    add_textbox(slide, verse_reference, Inches(8.5), Inches(0), Inches(1.5), Inches(0.75), PP_ALIGN.RIGHT, 18)
    add_textbox(slide, text_kor, Inches(0), Inches(0.5), Inches(20), Inches(20), PP_ALIGN.LEFT, 44)
    add_textbox(slide, text_eng, Inches(0), Inches(5.75), Inches(4), Inches(0.75), PP_ALIGN.LEFT, 18)

def generate_presentation_for_chapter(final_df):
    prs = initialize_presentation()
    for _, row in final_df.iterrows():
        # Convert NaN values to an empty string for both Korean and English text.
        text_kor = str(row['text_kor']) if pd.notna(row['text_kor']) else ""
        text_eng = str(row['text_eng']) if pd.notna(row['text_eng']) else ""

        # Now, pass these possibly modified strings to add_verse_slide.
        add_verse_slide(prs, row['book_name_kor'], row['book_name_eng'], row['chapter'], row['verse'], text_kor, text_eng)

    # Ensure that file_name is defined outside the loop to handle cases where final_df might be empty.
    if not final_df.empty:
        last_row = final_df.iloc[-1]
        file_name = f"{last_row['book_name_kor']}_{last_row['book_name_eng']}_{last_row['chapter']}.pptx"
        save_presentation(prs, file_name)
        print(f"Presentation saved as {file_name}")
    else:
        print("DataFrame is empty. No presentation was generated.")

def save_presentation(prs, file_name):
    prs.save(file_name)